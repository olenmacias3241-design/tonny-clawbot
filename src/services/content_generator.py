"""生成可直接下载的表格（CSV/Excel）和 PPT（.pptx）。

PPT 设计原则（参照常见优秀模版）：
- 层次清晰：封面主标题 > 页标题 > 正文，字号与字重区分
- 留白充足：边距与段落间距统一，每页要点不超过 6 条
- 配色统一：主色、强调色、正文色、辅助色成体系，与 Excel 表格风格一致
- 可选用参考模板：将优质 .pptx 放入 data/templates/pptx_template.pptx，生成时将沿用其版式与主题
"""

import csv
import io
import json
import re
import sys
from typing import List, Dict, Any, Optional

from src.bot.ai_provider import get_ai_provider
from src.utils.config import get_settings
from src.utils.logger import log

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    Presentation = Inches = Pt = RGBColor = MSO_SHAPE = PP_ALIGN = None


async def generate_table_data(prompt: str) -> List[Dict[str, str]]:
    """
    根据描述用 AI 生成表格数据，返回 list of dict（每行一个 dict，key 为列名）。
    """
    settings = get_settings()
    provider = get_ai_provider()
    system = """你是专业的表格数据生成助手。根据用户描述生成可直接使用的、高质量的表格数据。

要求：
1. 只输出一个合法的 JSON 数组，不要任何前后说明、不要 markdown 代码块、不要 ```。
2. 数组每个元素是一个对象，键为列名（中文或英文），值为该单元格的字符串或数字。同一列的类型必须一致（全是数字或全是字符串）。
3. 数据要真实、合理、多样：
   - 姓名/人名：使用多样化的中文姓名（如张三、李四、王芳、陈明等），不要重复或敷衍的“用户1、测试”。
   - 分数/成绩：在合理区间内（如 0–100），总分、平均分要按行正确计算。
   - 金额/数量：使用合理的整数或一位小数，单位与列名一致。
   - 日期：格式统一，如 2024-01-15 或 2024/1/15。
   - 部门/类别：与业务场景一致，如技术、产品、销售、市场等，多几类不要全一样。
4. 行数：一般 5–12 行；若用户明确要求“很多”或“大量”可生成 15–20 行。至少 5 行。
5. 列名清晰，与用户描述一致；若有“总分”“平均分”等，必须用数字填好并保证计算正确。

示例格式（仅作结构参考，请按用户描述生成实际内容）：
[{"姓名":"张三","部门":"技术","成绩":85},{"姓名":"李四","部门":"产品","成绩":92}]"""
    user = f"按下面描述生成表格，只输出一个 JSON 数组，不要其他任何文字：\n{prompt}"
    try:
        raw = await provider.generate_response(
            [{"role": "system", "content": system}, {"role": "user", "content": user}],
            temperature=0.3,
            max_tokens=4000,
        )
        raw = raw.strip()
        # 去掉可能的 markdown 包裹
        m = re.search(r"```(?:json)?\s*([\s\S]*?)```", raw)
        if m:
            raw = m.group(1).strip()
        # 尝试直接找 [...] 段
        arr = re.search(r"\[[\s\S]*\]", raw)
        if arr:
            raw = arr.group(0)
        data = json.loads(raw)
        if not isinstance(data, list) or not data:
            raise ValueError("AI 未返回有效表格数组")
        if not isinstance(data[0], dict):
            raise ValueError("表格每行应为对象")
        return data
    except json.JSONDecodeError as e:
        log.error(f"Table JSON parse error: {e}, raw: {raw[:200]}")
        raise ValueError("生成内容不是有效 JSON，请重试或换一种描述") from e


def table_to_csv(rows: List[Dict[str, str]]) -> str:
    """将 list of dict 转为 CSV 字符串（UTF-8，含 BOM 便于 Excel 识别）。"""
    if not rows:
        return ""
    out = io.StringIO()
    out.write("\ufeff")  # BOM for Excel
    writer = csv.DictWriter(out, fieldnames=list(rows[0].keys()), lineterminator="\n")
    writer.writeheader()
    writer.writerows(rows)
    return out.getvalue()


def table_to_xlsx(rows: List[Dict[str, str]]) -> bytes:
    """将 list of dict 转为带样式的 .xlsx（表头加粗/底色、边框、列宽、斑马纹）。"""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
        from openpyxl.utils import get_column_letter
    except ImportError:
        raise RuntimeError("请安装 openpyxl: pip install openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    if not rows:
        return io.BytesIO().getvalue()
    headers = list(rows[0].keys())
    # 表头样式：深灰底、白字、加粗、居中
    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    thin = Side(style="thin", color="B4B4B4")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    # 表头行
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.border = border
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    # 数据行 + 斑马纹
    light_fill = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")
    for r, row in enumerate(rows, 2):
        for c, h in enumerate(headers, 1):
            cell = ws.cell(row=r, column=c, value=row.get(h, ""))
            cell.border = border
            if (r - 2) % 2 == 1:
                cell.fill = light_fill
            cell.alignment = Alignment(vertical="center", wrap_text=True)
    # 冻结首行
    ws.freeze_panes = "A2"
    # 自动列宽（按内容，限制在 8–40 之间）
    for c, h in enumerate(headers, 1):
        col_letter = get_column_letter(c)
        max_w = len(str(h))
        for r in range(2, len(rows) + 2):
            val = rows[r - 2].get(h, "")
            max_w = min(40, max(max_w, len(str(val)) + 1))
        ws.column_dimensions[col_letter].width = max(8, max_w)
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


def _parse_structured_slides(text: str) -> Optional[Dict[str, Any]]:
    """
    解析「幻灯片1：封面」+ 标题/副标题/日期/姓名，以及「幻灯片N：标题」+ 要点列表。
    成功时返回 {"title", "subtitle", "date_author", "outline"}，否则返回 None。
    """
    if not text or "幻灯片" not in text:
        return None
    blocks = re.split(r"\n```|\n###\s*", text)
    cover_title = None
    cover_subtitle = None
    date_author_parts = []
    outline = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
        if not lines:
            continue
        first = lines[0]
        m = re.match(r"幻灯片\s*\d*\s*[：:]\s*(.+)", first)
        if not m:
            continue
        slide_label = m.group(1).strip()
        rest_lines = lines[1:]
        if "标题：" in block or "标题:" in block:
            for ln in rest_lines:
                if re.match(r"标题\s*[：:]\s*", ln):
                    cover_title = re.sub(r"标题\s*[：:]\s*", "", ln, count=1).strip()
                elif re.match(r"副标题\s*[：:]\s*", ln):
                    cover_subtitle = re.sub(r"副标题\s*[：:]\s*", "", ln, count=1).strip()
                elif re.match(r"日期\s*[：:]\s*", ln):
                    date_author_parts.append("日期：" + re.sub(r"日期\s*[：:]\s*", "", ln, count=1).strip())
                elif re.match(r"姓名\s*[：:]\s*", ln):
                    date_author_parts.append("姓名：" + re.sub(r"姓名\s*[：:]\s*", "", ln, count=1).strip())
            if cover_title and not outline:
                continue
        points = []
        for ln in rest_lines:
            if ln.startswith("- ") or ln.startswith("-"):
                points.append(ln[1:].strip().lstrip("- "))
            elif ln.startswith("• ") or ln.startswith("•"):
                points.append(ln[1:].strip().lstrip("• "))
        if slide_label and (points or ("标题：" not in block and "标题:" not in block)):
            if "封面" in slide_label and (cover_title or cover_subtitle) and not points:
                continue
            outline.append({"title": slide_label, "points": points})
    if not outline and cover_title is None:
        return None
    if cover_title is None and outline:
        cover_title = outline[0].get("title", "汇报")
    date_author = "  |  ".join(date_author_parts) if date_author_parts else None
    return {
        "title": cover_title,
        "subtitle": cover_subtitle or None,
        "date_author": date_author,
        "outline": outline,
    }


async def generate_ppt_outline(
    title: str,
    topic: str,
    conversation_context: Optional[str] = None,
) -> List[Dict[str, Any]]:
    """
    根据标题、主题与（可选）对话记录用 AI 生成 PPT 大纲（先出结构，有对话时尽量带关键信息），
    返回 [{"title":"幻灯片标题","points":["要点1","要点2"]}, ...]。
    若有 conversation_context，后续会再调用 expand 步骤把每页要点展开成具体内容。
    """
    provider = get_ai_provider()
    system = (
        "你是一个 PPT 大纲助手。根据用户给出的汇报标题、主题，以及可能提供的【对话记录】或【用户本次需求与资料】，生成 PPT 大纲。"
        "只输出一个 JSON 数组，不要任何其他文字。"
        "数组的每个元素：{\"title\":\"该页标题\",\"points\":[\"要点1\",\"要点2\",\"要点3\"]}。"
        "要求："
        "1. 若用户给的是完整需求说明（如包含「主要包括」「包括」「需要涵盖」或并列多条：批量xxx、批量xxx、传播到xxx），必须按需求中的**每一个要点单独成页或明确展开**，每页的 title 和 points 紧扣该要点，写出可执行、可落地的内容方向，不要笼统的「背景介绍」「总结」敷衍。"
        "2. 若有【对话记录】或【用户本次需求与资料】，必须从中提炼真实内容：各页标题和 points 要体现用户提到的具体模块、步骤、平台、工具，不要写空泛的模板句。"
        "3. 每页 3～6 条 points，每条是一句完整表述或关键信息（可先写简版，后续会再展开）。"
        "4. 通常 5～8 页：封面（标题页，points 可为空）、可选的目录、然后按用户需求的各模块逐页展开、最后总结与下一步。"
        "只输出 JSON 数组。"
    )
    user_parts = [f"汇报标题：{title}\n主题/内容方向：{topic}"]
    if conversation_context and conversation_context.strip():
        user_parts.append(f"\n【对话记录与资料】（请根据以下内容生成有实质内容的大纲）\n{conversation_context.strip()}")
    user_parts.append("\n请生成 PPT 大纲（JSON 数组）。")
    user = "".join(user_parts)
    try:
        raw = await provider.generate_response(
            [{"role": "system", "content": system}, {"role": "user", "content": user}],
            temperature=0.4,
            max_tokens=4000,
        )
        raw = raw.strip()
        m = re.search(r"```(?:json)?\s*([\s\S]*?)```", raw)
        if m:
            raw = m.group(1).strip()
        data = json.loads(raw)
        if not isinstance(data, list) or not data:
            raise ValueError("AI 未返回有效大纲")
        out = []
        for item in data:
            if isinstance(item, dict):
                out.append({
                    "title": str(item.get("title", "幻灯片")),
                    "points": [str(p) for p in item.get("points", []) if p],
                })
            else:
                out.append({"title": "幻灯片", "points": []})
        return out
    except json.JSONDecodeError as e:
        log.error(f"PPT outline JSON parse error: {e}, raw: {raw[:200]}")
        raise ValueError("生成内容不是有效 JSON，请重试") from e


async def expand_ppt_outline_with_conversation(
    outline: List[Dict[str, Any]],
    conversation_context: str,
    title: str,
) -> List[Dict[str, Any]]:
    """
    根据对话记录把大纲里每一页的 points 展开成具体、可读的要点（必须从对话中提炼数据、案例、结论），
    而不是只保留提纲。封面页不展开，其余页 3～6 条完整句。
    """
    if not (conversation_context and conversation_context.strip()):
        return outline
    provider = get_ai_provider()
    outline_json = json.dumps(outline, ensure_ascii=False, indent=2)
    system = (
        "你是 PPT 内容撰写助手。你会收到一份 PPT 大纲（JSON 数组）和【对话记录】或【用户本次需求与资料】。"
        "你的任务：把每一页的 points 扩展成 3～6 条**具体、可读、可执行**的要点，每条必须是一句或两句话。"
        "**硬性要求**："
        "1. 要点内容必须从【对话记录/用户需求与资料】中提炼，包含其中提到的具体模块、步骤、平台、工具、做法（如「批量创建文案」「批量生成视频」「传播到主流媒体」等），每一条都要有实质信息，禁止写「分析现状」「加强管理」等空泛句。"
        "2. 若用户需求里明确列了多条（如 A、B、C），对应页的 points 必须分别展开 A/B/C 的具体做法、工具或步骤，不要合并成笼统一句。"
        "封面页（通常是第一页）的 points 保持为空数组 []。"
        "只输出一个 JSON 数组，格式与输入大纲完全一致：每项 {\"title\":\"...\", \"points\":[\"...\",\"...\"]}，不要任何其他文字。"
    )
    user = f"汇报标题：{title}\n\n【当前大纲】\n{outline_json}\n\n【对话记录/用户本次需求与资料】（请根据以下内容填充每页的 points，写出具体、可执行的要点）\n{conversation_context.strip()}\n\n请输出扩展后的完整 JSON 数组。"
    try:
        raw = await provider.generate_response(
            [{"role": "system", "content": system}, {"role": "user", "content": user}],
            temperature=0.3,
            max_tokens=4000,
        )
        raw = raw.strip()
        m = re.search(r"```(?:json)?\s*([\s\S]*?)```", raw)
        if m:
            raw = m.group(1).strip()
        data = json.loads(raw)
        if not isinstance(data, list) or len(data) != len(outline):
            log.warning(f"Expand PPT outline returned different length, using original outline")
            return outline
        out = []
        for i, item in enumerate(data):
            if isinstance(item, dict):
                slide_title = str(item.get("title", outline[i].get("title", "幻灯片")))
                points = [str(p) for p in item.get("points", []) if p]
                out.append({"title": slide_title, "points": points[:8]})
            else:
                out.append(outline[i])
        return out
    except (json.JSONDecodeError, IndexError) as e:
        log.warning(f"Expand PPT outline failed: {e}, using original outline")
        return outline


def build_pptx_bytes(
    title: str,
    outline: List[Dict[str, Any]],
    subtitle: Optional[str] = None,
    date_author: Optional[str] = None,
) -> bytes:
    """根据标题和大纲生成 .pptx（内置版式：16:9、现代配色与排版，兼容且美观）。"""
    if Presentation is None:
        try:
            from pptx import Presentation as _P
            from pptx.util import Inches as _I, Pt as _Pt
            from pptx.dml.color import RGBColor as _R
            from pptx.enum.shapes import MSO_SHAPE as _M
            from pptx.enum.text import PP_ALIGN as _A
            globals()["Presentation"] = _P
            globals()["Inches"] = _I
            globals()["Pt"] = _Pt
            globals()["RGBColor"] = _R
            globals()["MSO_SHAPE"] = _M
            globals()["PP_ALIGN"] = _A
        except ImportError:
            raise RuntimeError(
                "未检测到 python-pptx。当前运行 Python："
                + sys.executable
                + " 。请在该环境中执行：pip install python-pptx，然后重启服务。"
            )

    # 16:9 幻灯片尺寸（更现代）
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    template_path = _get_pptx_template_path()
    if template_path:
        try:
            prs = Presentation(str(template_path))
            while len(prs.slides) > 0:
                i = len(prs.slides) - 1
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
            layouts = prs.slide_layouts
            title_lyt = layouts[0]
            content_lyt = layouts[1] if len(layouts) > 1 else layouts[0]
            s0 = prs.slides.add_slide(title_lyt)
            try:
                s0.shapes.title.text = title
            except Exception:
                pass
            if subtitle or date_author:
                sub_text = "\n".join([s for s in (subtitle, date_author) if s])
                try:
                    s0.placeholders[1].text = sub_text
                except (KeyError, IndexError):
                    pass
            for item in outline:
                s = prs.slides.add_slide(content_lyt)
                try:
                    s.shapes.title.text = item.get("title", "")
                except Exception:
                    pass
                points = item.get("points", [])[:8]
                try:
                    body = s.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    for i, point in enumerate(points):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = "●  " + (point if isinstance(point, str) else str(point))
                except (KeyError, IndexError):
                    pass
            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            return buf.read()
        except Exception as e:
            log.warning(f"PPT 参考模板加载失败，改用内置设计: {e}")

    # 现代配色：深色主色 + 单一强调色 + 充足留白
    primary = RGBColor(0x0F, 0x17, 0x2A)      # 深靛蓝 - 主色
    accent = RGBColor(0x0E, 0xA5, 0xE9)       # 天蓝 - 强调
    white = RGBColor(0xFF, 0xFF, 0xFF)
    body_color = RGBColor(0x33, 0x44, 0x55)   # 正文
    muted = RGBColor(0x64, 0x72, 0x85)        # 副标题
    card_bg = RGBColor(0xF8, 0xFA, 0xFC)      # 卡片背景
    card_border = RGBColor(0xE2, 0xE8, 0xF0)  # 细边框
    title_bar_bg = RGBColor(0x1E, 0x29, 0x3B) # 顶条

    FONT_TITLE = "Microsoft YaHei"
    FONT_BODY = "Microsoft YaHei"

    def _set_font(run, size_pt: int, bold: bool = False, color=None):
        try:
            run.font.name = FONT_BODY
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color
        except Exception:
            pass

    blank = prs.slide_layouts[6]

    def _is_section_slide(item: dict) -> bool:
        """无要点时用章节页版式（全幅色块 + 居中大标题）。"""
        points = item.get("points", []) or []
        return len(points) == 0

    def _is_key_message_slide(item: dict) -> bool:
        """1～2 条且较长时，用「要点/金句」版式（大号正文，无 bullet）。"""
        points = item.get("points", []) or []
        if len(points) not in (1, 2):
            return False
        avg_len = sum(len(str(p)) for p in points) / len(points)
        return avg_len >= 35

    # ---------- 封面：左侧竖条 + 大标题 + 副标题，右侧留白 ----------
    slide = prs.slides.add_slide(blank)
    # 左侧竖条（装饰）
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = accent
    left_bar.line.fill.background()
    # 标题区（居中偏左，留白多）
    tx = slide.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(10.8), Inches(1.6))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = primary
    try:
        p.font.name = FONT_TITLE
    except Exception:
        pass
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(16)
    if subtitle or date_author:
        sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.4), Inches(10.8), Inches(1.2))
        stf = sub_box.text_frame
        stf.word_wrap = True
        for i, line in enumerate([s for s in (subtitle, date_author) if s]):
            sp = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
            sp.text = line
            sp.font.size = Pt(16)
            sp.font.color.rgb = muted
            try:
                sp.font.name = FONT_BODY
            except Exception:
                pass
            sp.alignment = PP_ALIGN.LEFT
            sp.space_after = Pt(8)
    # 底部细线
    foot_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(7.0), Inches(2), Pt(2)
    )
    foot_line.fill.solid()
    foot_line.fill.fore_color.rgb = accent
    foot_line.line.fill.background()

    # ---------- 内容页：按类型选版式（章节页 / 要点页 / 普通内容页）----------
    for item in outline:
        slide = prs.slides.add_slide(blank)
        slide_title = item.get("title", "")
        points = item.get("points", [])[:8]

        if _is_section_slide(item):
            # 章节页：全幅色块 + 居中白字大标题
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.2), Inches(13.333), Inches(3.1)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.line.fill.background()
            tx = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.333), Inches(1.5))
            tf = tx.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title or " "
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = white
            try:
                p.font.name = FONT_TITLE
            except Exception:
                pass
            p.alignment = PP_ALIGN.CENTER
            continue

        if _is_key_message_slide(item):
            # 要点/金句页：页标题 + 大号正文（1～2 段），无 bullet
            strip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.14)
            )
            strip.fill.solid()
            strip.fill.fore_color.rgb = title_bar_bg
            strip.line.fill.background()
            vbar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Pt(4), Inches(0.65)
            )
            vbar.fill.solid()
            vbar.fill.fore_color.rgb = accent
            vbar.line.fill.background()
            title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.42), Inches(11.5), Inches(0.9))
            tft = title_box.text_frame
            tft.word_wrap = True
            pt = tft.paragraphs[0]
            pt.text = slide_title
            pt.font.size = Pt(26)
            pt.font.bold = True
            pt.font.color.rgb = primary
            try:
                pt.font.name = FONT_TITLE
            except Exception:
                pass
            body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5))
            tfb = body_box.text_frame
            tfb.word_wrap = True
            for i, point in enumerate(points):
                p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
                p.space_before = Pt(24)
                p.space_after = Pt(12)
                r = p.add_run()
                r.text = point if isinstance(point, str) else str(point)
                _set_font(r, 18, bold=False, color=body_color)
            continue

        # 普通内容页：顶条 + 左竖条 + 标题 + 卡片 + 要点
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.14)
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = title_bar_bg
        strip.line.fill.background()
        vbar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Pt(4), Inches(0.65)
        )
        vbar.fill.solid()
        vbar.fill.fore_color.rgb = accent
        vbar.line.fill.background()
        title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.42), Inches(11.5), Inches(0.9))
        tft = title_box.text_frame
        tft.word_wrap = True
        pt = tft.paragraphs[0]
        pt.text = slide_title
        pt.font.size = Pt(26)
        pt.font.bold = True
        pt.font.color.rgb = primary
        try:
            pt.font.name = FONT_TITLE
        except Exception:
            pass
        pt.space_after = Pt(4)
        content_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.35), Inches(12.2), Inches(5.75)
        )
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = card_bg
        content_bg.line.color.rgb = card_border
        content_bg.line.width = Pt(0.5)
        body_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.65), Inches(11.4), Inches(5.3))
        tfb = body_box.text_frame
        tfb.word_wrap = True
        for i, point in enumerate(points):
            p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
            p.space_before = Pt(18)
            p.space_after = Pt(8)
            r1 = p.add_run()
            r1.text = "  ◆  "
            _set_font(r1, 16, bold=False, color=accent)
            r2 = p.add_run()
            r2.text = point if isinstance(point, str) else str(point)
            _set_font(r2, 17, bold=False, color=body_color)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _generated_dir():
    from pathlib import Path
    root = Path(__file__).resolve().parent.parent.parent
    d = root / "data" / "generated"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _get_pptx_template_path():
    """若存在参考模板 data/templates/pptx_template.pptx 则返回其 Path，否则返回 None。"""
    from pathlib import Path
    root = Path(__file__).resolve().parent.parent.parent
    p = root / "data" / "templates" / "pptx_template.pptx"
    return p if p.is_file() else None


def _get_video_avatar_path():
    """若存在 data/templates/video_avatar.png 或 video_avatar.jpg 则返回 Path，用作视频讲解人物画面。"""
    from pathlib import Path
    root = Path(__file__).resolve().parent.parent.parent
    for name in ("video_avatar.png", "video_avatar.jpg"):
        p = root / "data" / "templates" / name
        if p.is_file():
            return p
    return None


def _get_sadtalker_dir():
    """若已安装 SadTalker（含 inference.py 与 checkpoints），返回其根目录 Path，否则返回 None。"""
    import os
    from pathlib import Path
    root = Path(__file__).resolve().parent.parent.parent
    candidates = [
        os.environ.get("SADTALKER_DIR"),
        root / "SadTalker",
        root / "sadtalker",
    ]
    for c in candidates:
        if not c:
            continue
        d = Path(c).resolve()
        if not d.is_dir():
            continue
        if (d / "inference.py").is_file() and (d / "checkpoints").is_dir():
            return d
    return None


def _run_sadtalker(sadtalker_dir, source_image, driven_audio, result_dir, timeout=480):
    """
    在 sadtalker_dir 下执行 inference.py，用 source_image + driven_audio 生成口播视频。
    返回 (生成的 .mp4 的 Path 或 None, 失败时的 stderr 摘要)。
    若设置 SADTALKER_PYTHON，则用该 Python 执行（可指向含正确 torch/torchvision 的 conda/venv）。
    """
    import subprocess
    import os
    import sys
    from pathlib import Path
    from src.utils.config import get_settings
    py = os.environ.get("SADTALKER_PYTHON") or get_settings().sadtalker_python or sys.executable
    sadtalker_dir = Path(sadtalker_dir).resolve()
    result_dir = Path(result_dir).resolve()
    result_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        py, str(sadtalker_dir / "inference.py"),
        "--driven_audio", str(Path(driven_audio).resolve()),
        "--source_image", str(Path(source_image).resolve()),
        "--result_dir", str(result_dir),
        "--preprocess", "full",
        "--still",
    ]
    try:
        r = subprocess.run(
            cmd, cwd=str(sadtalker_dir), capture_output=True, timeout=timeout,
            text=True, encoding="utf-8", errors="replace",
        )
        if r.returncode != 0:
            stderr_snippet = (r.stderr or "")[-800:] if (r.stderr or "") else ""
            log.warning(f"SadTalker exit {r.returncode}: {stderr_snippet}")
            return None, stderr_snippet
        found = list(result_dir.glob("*.mp4"))
        return (Path(found[0]) if found else None, "")
    except (FileNotFoundError, subprocess.TimeoutExpired) as e:
        err_msg = str(e)
        log.warning(f"SadTalker run failed: {e}")
        return None, err_msg


async def generate_table_and_save(prompt: str):
    """生成表格并保存为 CSV 和 Excel（若已安装 openpyxl），返回 (csv_name, xlsx_name 或 None)。"""
    import uuid
    log.info("generate_table_and_save: calling AI for table data")
    rows = await generate_table_data(prompt)
    log.info(f"generate_table_and_save: got {len(rows)} rows")
    uid = uuid.uuid4().hex[:12]
    base = _generated_dir()
    csv_name = f"table_{uid}.csv"
    csv_path = base / csv_name
    csv_path.write_text(table_to_csv(rows), encoding="utf-8-sig")
    log.info(f"generate_table_and_save: wrote {csv_path}")
    xlsx_name = None
    try:
        xlsx_name = f"table_{uid}.xlsx"
        (base / xlsx_name).write_bytes(table_to_xlsx(rows))
        log.info(f"generate_table_and_save: wrote xlsx {xlsx_name}")
    except Exception as e:
        log.warning(f"Excel 生成跳过（可安装 openpyxl 后重试）: {e}")
    return csv_name, xlsx_name


async def generate_document_content(prompt: str, conversation_context: Optional[str] = None) -> str:
    """根据用户描述与（可选）对话/资料用 AI 生成一篇完整文档正文（Markdown），内容必须基于对话与资料，不是提纲或模板。"""
    provider = get_ai_provider()
    system = """你是文档撰写助手。根据用户描述以及可能提供的【对话记录】或用户粘贴的【资料】，撰写一篇**完整、可直接使用**的 Markdown 文档正文（不是提纲也不是模板）。
硬性要求：
- 若有【对话记录】或用户提供的资料，文档内容必须**基于其中的具体信息**撰写：使用对话里的数据、案例、结论、人名、产品名等，写出完整段落和列表，禁止只写小标题或空泛句。
- 文档应是成形的正文：有 # ## ### 标题、有段落、有列表或要点，总长度一般 500～2000 字（按需求调整），读者可直接当正式文档使用。
- 使用 Markdown 语法，段落间空行，只输出文档正文，不要「以下是文档」等前言。"""
    user_parts = [f"用户本次需求：\n{prompt}"]
    if conversation_context and conversation_context.strip():
        user_parts.append(f"\n\n【对话记录与用户提供的资料】（请根据以下内容写出完整文档，不要只给提纲）\n{conversation_context.strip()}")
    user = "".join(user_parts)
    raw = await provider.generate_response(
        [{"role": "system", "content": system}, {"role": "user", "content": user}],
        temperature=0.4,
        max_tokens=6000,
    )
    return (raw or "").strip()


async def generate_document_and_save(prompt: str, conversation_context: Optional[str] = None) -> str:
    """生成文档并保存为 .md，返回文件名。可传入对话上下文以基于对话内容生成。"""
    import uuid
    log.info("generate_document_and_save: calling AI for document content")
    content = await generate_document_content(prompt, conversation_context=conversation_context)
    uid = uuid.uuid4().hex[:12]
    base = _generated_dir()
    name = f"doc_{uid}.md"
    base.joinpath(name).write_text(content, encoding="utf-8")
    log.info(f"generate_document_and_save: wrote {base / name}")
    return name


def _markdown_to_docx_bytes(md_text: str) -> bytes:
    """将 Markdown 文本转为 .docx 字节（标题、段落、列表），中文字体+行距避免乱码与显示不全。"""
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.enum.text import WD_LINE_SPACING
    except ImportError:
        raise RuntimeError("请安装 python-docx: pip install python-docx")
    doc = Document()
    # 中英文统一用同一字体，避免混排时行高不一致导致中文被裁切；宋体在 Windows/macOS 较常见
    cjk_font = "SimSun"  # 宋体

    def set_run_cjk_font(run):
        """设置 run 的 ascii/hAnsi/eastAsia 为同一中文字体，避免显示不全或乱码。"""
        try:
            run.font.name = cjk_font
            rPr = run._element.get_or_add_rPr()
            rFonts = rPr.get_or_add_rFonts()
            rFonts.set(qn("w:eastAsia"), cjk_font)
            rFonts.set(qn("w:ascii"), cjk_font)
            rFonts.set(qn("w:hAnsi"), cjk_font)
        except Exception:
            pass

    def set_style_cjk_font(style):
        try:
            style.font.name = cjk_font
            if style._element.rPr is not None:
                rFonts = style._element.rPr.get_or_add_rFonts()
                rFonts.set(qn("w:eastAsia"), cjk_font)
                rFonts.set(qn("w:ascii"), cjk_font)
                rFonts.set(qn("w:hAnsi"), cjk_font)
        except Exception:
            pass

    def set_paragraph_line_spacing(paragraph, multiplier=1.15):
        """设置段落行距，避免中文上下被裁切。"""
        try:
            paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
            paragraph.paragraph_format.line_spacing = multiplier
        except Exception:
            pass

    for style_name in ("Normal", "Heading 1", "Heading 2", "Heading 3"):
        if style_name in doc.styles:
            set_style_cjk_font(doc.styles[style_name])

    lines = md_text.split("\n")
    i = 0
    in_ul = False
    in_ol = False
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            if in_ul or in_ol:
                in_ul = in_ol = False
            i += 1
            continue
        if stripped.startswith("###"):
            if in_ul or in_ol:
                in_ul = in_ol = False
            doc.add_heading(stripped.lstrip("#").strip(), level=3)
            para = doc.paragraphs[-1]
            for run in para.runs:
                set_run_cjk_font(run)
            set_paragraph_line_spacing(para, 1.2)
            i += 1
            continue
        if stripped.startswith("##"):
            if in_ul or in_ol:
                in_ul = in_ol = False
            doc.add_heading(stripped.lstrip("#").strip(), level=2)
            para = doc.paragraphs[-1]
            for run in para.runs:
                set_run_cjk_font(run)
            set_paragraph_line_spacing(para, 1.2)
            i += 1
            continue
        if stripped.startswith("#"):
            if in_ul or in_ol:
                in_ul = in_ol = False
            doc.add_heading(stripped.lstrip("#").strip(), level=1)
            para = doc.paragraphs[-1]
            for run in para.runs:
                set_run_cjk_font(run)
            set_paragraph_line_spacing(para, 1.2)
            i += 1
            continue
        if re.match(r"^[\-\*]\s+", stripped) or stripped.startswith("- ") or stripped.startswith("* "):
            text = re.sub(r"^[\-\*]\s+", "", stripped, count=1)
            p = doc.add_paragraph(style="List Bullet")
            r = p.add_run(text)
            set_run_cjk_font(r)
            set_paragraph_line_spacing(p, 1.15)
            in_ul = True
            in_ol = False
            i += 1
            continue
        if re.match(r"^\d+\.\s+", stripped):
            text = re.sub(r"^\d+\.\s+", "", stripped, count=1)
            p = doc.add_paragraph(style="List Number")
            r = p.add_run(text)
            set_run_cjk_font(r)
            set_paragraph_line_spacing(p, 1.15)
            in_ol = True
            in_ul = False
            i += 1
            continue
        in_ul = in_ol = False
        p = doc.add_paragraph(stripped)
        for run in p.runs:
            set_run_cjk_font(run)
        set_paragraph_line_spacing(p, 1.15)
        i += 1
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


async def generate_docx_and_save(prompt: str, conversation_context: Optional[str] = None) -> str:
    """根据用户描述与（可选）对话记录用 AI 生成文档并保存为 .docx，返回文件名。"""
    import uuid
    log.info("generate_docx_and_save: calling AI for document content")
    content = await generate_document_content(prompt, conversation_context=conversation_context)
    docx_bytes = _markdown_to_docx_bytes(content)
    uid = uuid.uuid4().hex[:12]
    base = _generated_dir()
    name = f"word_{uid}.docx"
    (base / name).write_bytes(docx_bytes)
    log.info(f"generate_docx_and_save: wrote {base / name}")
    return name


async def generate_ppt_and_save(
    title: str,
    topic: str,
    conversation_context: Optional[str] = None,
) -> str:
    """根据标题、主题与（可选）对话记录用 AI 生成大纲，有对话时再展开成具体内容，然后保存 .pptx。"""
    import uuid
    outline = await generate_ppt_outline(title, topic, conversation_context=conversation_context)
    if conversation_context and conversation_context.strip():
        outline = await expand_ppt_outline_with_conversation(outline, conversation_context, title)
    content = build_pptx_bytes(title, outline)
    uid = uuid.uuid4().hex[:12]
    base = _generated_dir()
    name = f"ppt_{uid}.pptx"
    (base / name).write_bytes(content)
    return name


def generate_ppt_from_structured_text(text: str) -> Optional[str]:
    """
    若文本为「幻灯片1：封面」+ 标题/副标题/日期/姓名 及 幻灯片N + 要点 的格式，则解析并生成 .pptx，返回文件名；否则返回 None。
    """
    parsed = _parse_structured_slides(text)
    if not parsed or not parsed.get("outline"):
        return None
    import uuid
    content = build_pptx_bytes(
        parsed["title"],
        parsed["outline"],
        subtitle=parsed.get("subtitle"),
        date_author=parsed.get("date_author"),
    )
    uid = uuid.uuid4().hex[:12]
    base = _generated_dir()
    name = f"ppt_{uid}.pptx"
    (base / name).write_bytes(content)
    return name


async def generate_video_script(
    conversation_context: str,
    user_prompt: str = "",
) -> Dict[str, Any]:
    """根据对话内容用 AI 生成视频脚本与分镜，保存到 data/generated。"""
    import uuid
    provider = get_ai_provider()
    system = """你是短视频/口播脚本与分镜师。根据用户提供的【对话记录】生成一份可直接用于拍摄或配音的视频脚本。
输出一个 JSON 对象，包含：
- "script": 完整旁白稿（适合朗读配音），基于对话内容提炼，口语化。
- "shots": 分镜数组，每项 {"scene": 序号, "duration_sec": 建议秒数, "visual": "画面描述", "narration": "该镜头旁白"}，5～15 个镜头。
内容必须来自对话中的具体信息。只输出一个 JSON，不要 markdown 包裹外的文字。"""
    user = "根据以下对话内容生成视频脚本与分镜。\n"
    if user_prompt and user_prompt.strip():
        user += f"用户补充要求：{user_prompt.strip()}\n\n"
    user += f"【对话记录】\n{(conversation_context or '').strip() or '（无）'}\n\n请输出 JSON。"
    try:
        raw = await provider.generate_response(
            [{"role": "system", "content": system}, {"role": "user", "content": user}],
            temperature=0.4,
            max_tokens=4000,
        )
        raw = (raw or "").strip()
        m = re.search(r"```(?:json)?\s*([\s\S]*?)```", raw)
        if m:
            raw = m.group(1).strip()
        obj = json.loads(raw)
        script = str(obj.get("script", ""))
        shots = obj.get("shots") if isinstance(obj.get("shots"), list) else []
        uid = uuid.uuid4().hex[:12]
        base = _generated_dir()
        script_name = f"video_script_{uid}.txt"
        (base / script_name).write_text(script, encoding="utf-8")
        json_name = f"video_script_{uid}.json"
        (base / json_name).write_text(
            json.dumps({"script": script, "shots": shots}, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        log.info(f"generate_video_script: wrote {base / script_name}")
        return {"script": script, "shots": shots, "script_filename": script_name, "json_filename": json_name}
    except json.JSONDecodeError as e:
        log.error(f"Video script JSON parse error: {e}")
        raise ValueError("生成内容不是有效 JSON，请重试") from e


async def generate_video_file_from_script(script: str) -> Dict[str, Any]:
    """
    根据旁白稿用 TTS 生成语音，再与静态画面合成为 .mp4。
    返回 {"video_filename": str|None, "audio_filename": str|None, "error": str|None}。
    若未安装 ffmpeg 但 TTS 成功，会保留并返回 audio_filename 供下载配音。
    """
    import uuid
    import subprocess
    import shutil
    result = {"video_filename": None, "audio_filename": None, "error": None}
    script = (script or "").strip()
    if not script:
        result["error"] = "旁白稿为空"
        return result
    text_for_tts = script[:3500] if len(script) > 3500 else script
    base = _generated_dir()
    uid = uuid.uuid4().hex[:12]
    mp3_path = base / f"video_audio_{uid}.mp3"
    mp4_name = f"video_{uid}.mp4"
    mp4_path = base / mp4_name
    try:
        import edge_tts
        comm = edge_tts.Communicate(text_for_tts, voice="zh-CN-YunxiNeural")
        await comm.save(str(mp3_path))
        if not mp3_path.is_file():
            result["error"] = "语音合成失败（edge-tts 未生成文件）"
            return result
    except Exception as e:
        err = str(e)
        if "nodename nor servname" in err or "Cannot connect" in err:
            result["error"] = "语音合成失败：无法连接 TTS 服务，请检查网络或代理。"
        else:
            result["error"] = "语音合成失败（edge-tts）：" + err[:200]
        log.warning(f"edge-tts failed: {e}")
        return result
    ffmpeg = shutil.which("ffmpeg")
    if not ffmpeg:
        result["audio_filename"] = mp3_path.name
        result["error"] = "未检测到 ffmpeg，无法生成 MP4。已生成配音 MP3，可先下载收听；安装 ffmpeg 后可生成视频。"
        log.info(f"generate_video_file: no ffmpeg, kept audio {mp3_path}")
        return result
    # 获取音频时长
    try:
        out = subprocess.run(
            [shutil.which("ffprobe") or "ffprobe", "-v", "error", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", str(mp3_path)],
            capture_output=True, text=True, timeout=10
        )
        duration_sec = float((out.stdout or "0").strip() or 0)
    except Exception:
        duration_sec = 60.0
    if duration_sec <= 0:
        duration_sec = 60.0
    # 按句拆分旁白，生成 SRT 字幕
    chunks = re.split(r"(?<=[。！？；])\s*|(?<=[\n])", text_for_tts)
    chunks = [c.strip() for c in chunks if c.strip()][:50]
    if not chunks:
        chunks = [text_for_tts[:100] + ("…" if len(text_for_tts) > 100 else "")]
    n = len(chunks)
    chunk_dur = duration_sec / n
    srt_parts = []
    for i, line in enumerate(chunks):
        start = i * chunk_dur
        end = (i + 1) * chunk_dur
        srt_parts.append(f"{i+1}\n{_srt_ts(start)} --> {_srt_ts(end)}\n{line}\n")
    srt_path = base / f"video_srt_{uid}.srt"
    srt_path.write_text("\n".join(srt_parts), encoding="utf-8")

    # 使用 /tmp 下无冒号的路径给 subtitles 滤镜，避免 filter 解析时把路径里的 : 当成分隔符
    import tempfile
    from pathlib import Path
    srt_for_ffmpeg = Path(tempfile.gettempdir()) / f"claw_srt_{uid}.srt"
    try:
        srt_for_ffmpeg.write_text("\n".join(srt_parts), encoding="utf-8")
    except Exception:
        srt_for_ffmpeg = srt_path
    # 当前用软字幕混流（不依赖 libass）；烧录字幕需 ffmpeg 带 libass

    used_sadtalker = False
    avatar_path = _get_video_avatar_path()
    sadtalker_dir = _get_sadtalker_dir()
    sadtalker_fail_msg = ""
    try:
        if avatar_path and avatar_path.is_file() and sadtalker_dir:
            wav_path = base / f"sadtalker_audio_{uid}.wav"
            try:
                subprocess.run(
                    [ffmpeg, "-y", "-i", str(mp3_path), "-acodec", "pcm_s16le", "-ar", "16000", str(wav_path)],
                    capture_output=True, timeout=30, check=True,
                )
            except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
                wav_path = None
            if wav_path and wav_path.is_file():
                result_dir = base / f"sadtalker_out_{uid}"
                out_mp4, sadtalker_fail_msg = _run_sadtalker(sadtalker_dir, avatar_path, wav_path, result_dir)
                try:
                    wav_path.unlink(missing_ok=True)
                except Exception:
                    pass
                if out_mp4 and out_mp4.is_file():
                    try:
                        subprocess.run(
                            [
                                ffmpeg, "-y", "-i", str(out_mp4), "-i", str(srt_for_ffmpeg.resolve()),
                                "-c:v", "copy", "-c:a", "copy", "-c:s", "mov_text",
                                "-metadata:s:s:0", "language=chi", str(mp4_path.resolve()),
                            ],
                            capture_output=True, timeout=120, check=True,
                        )
                        used_sadtalker = True
                    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
                        pass
                    try:
                        shutil.rmtree(result_dir, ignore_errors=True)
                    except Exception:
                        pass

        if not used_sadtalker:
            avatar_path = avatar_path or _get_video_avatar_path()
            # 两段式合成：先生成无字幕视频，再单独烧录字幕，避免单条 filter 链+输出导致 Invalid argument
            mp4_naked = base / f"video_{uid}_naked.mp4"
            try:
                if avatar_path and avatar_path.is_file():
                    cmd1 = [
                        ffmpeg, "-y", "-loop", "1", "-i", str(avatar_path.resolve()),
                        "-i", str(mp3_path), "-vf", "scale=1280:720:force_original_aspect_ratio=decrease,pad=1280:720:(ow-iw)/2:(oh-ih)/2:black",
                        "-c:v", "libx264", "-tune", "stillimage", "-c:a", "aac",
                        "-shortest", "-pix_fmt", "yuv420p", str(mp4_naked.resolve()),
                    ]
                else:
                    cmd1 = [
                        ffmpeg, "-y", "-f", "lavfi", "-i", "color=c=#1e1e24:s=1280x720:d=600",
                        "-i", str(mp3_path), "-c:v", "libx264", "-tune", "stillimage", "-c:a", "aac",
                        "-shortest", "-pix_fmt", "yuv420p", str(mp4_naked.resolve()),
                    ]
                r1 = subprocess.run(cmd1, capture_output=True, timeout=120, text=True, encoding="utf-8", errors="replace")
                if r1.returncode != 0:
                    stderr_full = (r1.stderr or "").strip()
                    stderr_msg = stderr_full[-500:] if len(stderr_full) > 500 else stderr_full
                    raise RuntimeError(stderr_msg or f"ffmpeg 第一步退出码 {r1.returncode}")
                if not mp4_naked.is_file():
                    raise RuntimeError("ffmpeg 未生成无字幕视频文件")
                # 先尝试烧录字幕（需 ffmpeg 带 libass）；失败则用软字幕
                filter_script = Path(tempfile.gettempdir()) / f"claw_filter_{uid}.txt"
                filter_line = f"[0:v]subtitles='{srt_for_ffmpeg.resolve()}'[v]\n"
                try:
                    filter_script.write_text(filter_line, encoding="utf-8")
                except Exception:
                    filter_script = None
                step2_ok = False
                if filter_script and filter_script.is_file():
                    cmd2_burn = [
                        ffmpeg, "-y", "-i", str(mp4_naked.resolve()),
                        "-filter_complex_script", str(filter_script.resolve()),
                        "-map", "[v]", "-map", "0:a", "-c:a", "copy", "-c:v", "libx264", "-pix_fmt", "yuv420p",
                        str(mp4_path.resolve()),
                    ]
                    r2 = subprocess.run(cmd2_burn, capture_output=True, timeout=120, text=True, encoding="utf-8", errors="replace")
                    if r2.returncode == 0:
                        step2_ok = True
                    elif "No such filter" in (r2.stderr or "") or "Filter not found" in (r2.stderr or ""):
                        pass  # 无 libass，下面用软字幕
                try:
                    if filter_script:
                        filter_script.unlink(missing_ok=True)
                except Exception:
                    pass
                if not step2_ok:
                    cmd2 = [
                        ffmpeg, "-y", "-i", str(mp4_naked.resolve()), "-i", str(srt_for_ffmpeg.resolve()),
                        "-c:v", "copy", "-c:a", "copy", "-c:s", "mov_text", "-metadata:s:s:0", "language=chi",
                        str(mp4_path.resolve()),
                    ]
                    r2 = subprocess.run(cmd2, capture_output=True, timeout=120, text=True, encoding="utf-8", errors="replace")
                    if r2.returncode != 0:
                        stderr_full = (r2.stderr or "").strip()
                        stderr_msg = stderr_full[-500:] if len(stderr_full) > 500 else stderr_full
                        raise RuntimeError(stderr_msg or f"ffmpeg 第二步(字幕)退出码 {r2.returncode}")
            finally:
                try:
                    mp4_naked.unlink(missing_ok=True)
                except Exception:
                    pass
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired, RuntimeError) as e:
        result["audio_filename"] = mp3_path.name
        err_detail = str(e)
        if hasattr(e, "stderr") and e.stderr:
            err_detail = (err_detail + " " + (e.stderr[:300] if isinstance(e.stderr, str) else e.stderr.decode("utf-8", errors="replace")[:300])).strip()
        result["error"] = "视频合成失败（ffmpeg/SadTalker），已保留配音 MP3 可下载。错误：" + err_detail[:500]
        log.warning(f"video synthesis failed: {e}")
        return result
    finally:
        if mp4_path.is_file():
            try:
                mp3_path.unlink()
            except Exception:
                pass
        try:
            srt_path.unlink(missing_ok=True)
        except Exception:
            pass
        if srt_for_ffmpeg != srt_path:
            try:
                srt_for_ffmpeg.unlink(missing_ok=True)
            except Exception:
                pass
    if mp4_path.is_file():
        log.info(f"generate_video_file: wrote {mp4_path}")
        result["video_filename"] = mp4_name
        result["video_sadtalker_used"] = used_sadtalker
        if not used_sadtalker and (sadtalker_dir or sadtalker_fail_msg):
            note = "本次为静态人物图，未使用口播动嘴。"
            if sadtalker_fail_msg:
                snippet = sadtalker_fail_msg.replace("\n", " ").strip()[:220]
                note += " SadTalker 报错：" + snippet
                if "Can't get the coeffs" in sadtalker_fail_msg or "first_coeff_path" in sadtalker_fail_msg or "3DMM" in sadtalker_fail_msg:
                    note += "（动漫/卡通图常无法通过人脸检测，可换真人照片试口播，或继续使用静态图。）"
                elif "timed out" in sadtalker_fail_msg or "TimeoutExpired" in sadtalker_fail_msg:
                    note += "（口播生成超时 8 分钟，已改为静态图。可缩短旁白或稍后重试。）"
            else:
                note += " 若需口播，请在 SadTalker 目录执行 pip install -r requirements.txt 后重试。"
            result["video_note"] = note
    return result


def _srt_ts(sec: float) -> str:
    h = int(sec // 3600)
    m = int((sec % 3600) // 60)
    s = sec % 60
    return f"{h:02d}:{m:02d}:{int(s):02d},{int((s % 1) * 1000):03d}"
